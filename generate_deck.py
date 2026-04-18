"""
Generate a 10-slide realistic startup pitch deck for NovaMed AI.
Uses python-pptx to create a polished PPTX file.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Palette ──────────────────────────────────────────────────
DARK_BG     = RGBColor(0x0F, 0x17, 0x2A)   # Deep navy
ACCENT_BLUE = RGBColor(0x00, 0x9E, 0xF7)   # Electric blue
ACCENT_TEAL = RGBColor(0x00, 0xD4, 0xAA)   # Teal-green
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xB0, 0xB8, 0xC8)
ORANGE      = RGBColor(0xFF, 0x8C, 0x42)
MEDIUM_GRAY = RGBColor(0x1A, 0x25, 0x3C)
CARD_BG     = RGBColor(0x16, 0x20, 0x37)

OUTPUT_PATH = os.path.join("data", "raw_decks", "NovaMed_AI_Pitch_Deck.pptx")


def set_slide_background(slide, color):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Helper to add a styled text box."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, paragraphs_data,
                   font_name="Calibri"):
    """Add a text box with multiple styled paragraphs.
    paragraphs_data: list of dicts with keys: text, size, color, bold, alignment, space_after
    """
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                      Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pdata in enumerate(paragraphs_data):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = pdata.get("text", "")
        p.font.size = Pt(pdata.get("size", 14))
        p.font.color.rgb = pdata.get("color", WHITE)
        p.font.bold = pdata.get("bold", False)
        p.font.name = font_name
        p.alignment = pdata.get("alignment", PP_ALIGN.LEFT)
        if "space_after" in pdata:
            p.space_after = Pt(pdata["space_after"])
        if "space_before" in pdata:
            p.space_before = Pt(pdata["space_before"])
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_accent_line(slide, left, top, width, color=ACCENT_BLUE):
    """Add a thin accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(0.04)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Slide 1 – Title / Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_background(slide, DARK_BG)

    add_accent_line(slide, 1.5, 1.8, 2.0, ACCENT_TEAL)

    add_textbox(slide, 1.5, 2.0, 7.0, 1.0, "NovaMed AI", 44,
                ACCENT_BLUE, bold=True, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 1.5, 2.8, 7.0, 0.8,
                "AI-Powered Diagnostic Intelligence for Rural Healthcare",
                20, ACCENT_TEAL, bold=False, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 1.5, 3.8, 7.0, 0.5,
                "Series A  •  $12M Round  •  April 2026",
                14, LIGHT_GRAY, alignment=PP_ALIGN.LEFT)

    add_multi_text(slide, 1.5, 5.0, 7.0, 1.5, [
        {"text": "Dr. Ananya Sharma, CEO & Co-Founder", "size": 12, "color": WHITE, "bold": True, "space_after": 4},
        {"text": "ananya@novamed.ai  |  +91-98765-43210", "size": 10, "color": LIGHT_GRAY, "space_after": 2},
        {"text": "www.novamed.ai", "size": 10, "color": ACCENT_BLUE, "bold": True},
    ])


def slide_02_problem(prs):
    """Slide 2 – Problem Statement"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "THE PROBLEM", 28,
                ACCENT_BLUE, bold=True)
    add_accent_line(slide, 0.8, 1.05, 1.2, ACCENT_BLUE)

    # Stat cards
    stats = [
        ("68%", "of India's population lives in\nrural areas with < 1 doctor\nper 10,000 people"),
        ("4.2 hrs", "average travel time for a\nrural patient to reach the\nnearest diagnostic center"),
        ("₹18,000 Cr", "lost annually in GDP due\nto delayed diagnosis and\npreventable disease escalation"),
    ]
    for i, (big, desc) in enumerate(stats):
        x = 0.8 + i * 2.9
        add_rounded_rect(slide, x, 1.4, 2.6, 2.2)
        add_textbox(slide, x + 0.2, 1.6, 2.2, 0.6, big, 26,
                    ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.2, 2.3, 2.2, 1.2, desc, 11,
                    LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_multi_text(slide, 0.8, 3.9, 8.0, 2.5, [
        {"text": "India's rural healthcare crisis is not a supply problem — it's an intelligence gap.",
         "size": 16, "color": WHITE, "bold": True, "space_after": 12},
        {"text": "• PHCs (Primary Health Centers) have equipment but lack trained radiologists and pathologists.",
         "size": 11, "color": LIGHT_GRAY, "space_after": 6},
        {"text": "• 73% of misdiagnoses in Tier-3 towns originate from human fatigue during high patient loads (>80/day).",
         "size": 11, "color": LIGHT_GRAY, "space_after": 6},
        {"text": "• Government schemes (Ayushman Bharat) cover treatment costs — but patients arrive too late because early diagnosis never happened.",
         "size": 11, "color": LIGHT_GRAY},
    ])


def slide_03_solution(prs):
    """Slide 3 – Solution"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "OUR SOLUTION", 28,
                ACCENT_TEAL, bold=True)
    add_accent_line(slide, 0.8, 1.05, 1.2, ACCENT_TEAL)

    add_textbox(slide, 0.8, 1.3, 8.0, 0.8,
                "NovaMed AI is an edge-deployed diagnostic assistant that brings specialist-grade "
                "analysis to every Primary Health Center in India.",
                14, WHITE, bold=False)

    features = [
        ("🔬  AI Radiology Engine",
         "Analyzes X-ray, ultrasound & retinal scans with 96.3% accuracy "
         "(validated on 1.2M anonymized images from AIIMS & CMC Vellore). "
         "Runs on-device in <8 seconds per scan."),
        ("🩸  Smart Lab Interpreter",
         "Processes CBC, LFT, KFT reports and flags critical deviations "
         "with differential diagnosis suggestions. Integrates with 14 common "
         "Indian lab-machine OEMs via HL7/FHIR."),
        ("📱  Offline-First Architecture",
         "Full functionality on ₹12,000 Android tablets with no internet. "
         "Data syncs automatically when connectivity resumes. "
         "Designed for 2G/3G networks."),
        ("🗣️  Multilingual Voice UX",
         "Hindi, Tamil, Telugu, Bengali, Marathi voice-guided workflows so "
         "ASHA workers with minimal training can operate the system independently."),
    ]

    for i, (title, desc) in enumerate(features):
        y = 2.3 + i * 1.15
        add_rounded_rect(slide, 0.8, y, 8.2, 1.0)
        add_textbox(slide, 1.0, y + 0.08, 7.8, 0.35, title, 13,
                    ACCENT_BLUE, bold=True)
        add_textbox(slide, 1.0, y + 0.45, 7.8, 0.55, desc, 10,
                    LIGHT_GRAY)


def slide_04_market(prs):
    """Slide 4 – Market Opportunity (TAM/SAM/SOM)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "MARKET OPPORTUNITY", 28,
                ACCENT_BLUE, bold=True)
    add_accent_line(slide, 0.8, 1.05, 1.2, ACCENT_BLUE)

    # TAM / SAM / SOM cards
    market_data = [
        ("TAM", "$47B", "Global AI in Healthcare\nDiagnostics (2028E)\nCAGR 38.4% — Grand View Research"),
        ("SAM", "$8.2B", "India AI-Assisted Diagnostics\n+ Rural Health-Tech Market\n(2028E)"),
        ("SOM", "$420M", "Year 5 Revenue Target\n155,000 PHCs × ₹28,000\nSaaS Annual License"),
    ]
    for i, (label, value, desc) in enumerate(market_data):
        x = 0.8 + i * 2.9
        add_rounded_rect(slide, x, 1.4, 2.6, 2.8)
        add_textbox(slide, x, 1.55, 2.6, 0.4, label, 14,
                    ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, 1.95, 2.6, 0.55, value, 30,
                    WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.15, 2.6, 2.3, 1.4, desc, 10,
                    LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_multi_text(slide, 0.8, 4.5, 8.0, 2.0, [
        {"text": "Why Now?", "size": 16, "color": ACCENT_TEAL, "bold": True, "space_after": 8},
        {"text": "• National Digital Health Mission (ABDM) mandates digitized records at all PHCs by 2027 — creating a greenfield integration opportunity.",
         "size": 11, "color": LIGHT_GRAY, "space_after": 5},
        {"text": "• Edge AI chips (Qualcomm QCS6490) now deliver 12 TOPS in a ₹3,000 module — making on-device inference commercially viable for the first time.",
         "size": 11, "color": LIGHT_GRAY, "space_after": 5},
        {"text": "• 62% of state NHM budgets now have a dedicated line-item for health-tech procurement (up from 18% in 2022).",
         "size": 11, "color": LIGHT_GRAY},
    ])


def slide_05_traction(prs):
    """Slide 5 – Traction & Milestones"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "TRACTION & MILESTONES", 28,
                ACCENT_TEAL, bold=True)
    add_accent_line(slide, 0.8, 1.05, 1.2, ACCENT_TEAL)

    # Key metric cards (row 1)
    metrics_row1 = [
        ("312", "PHCs Live", "across Rajasthan, UP,\nMadhya Pradesh"),
        ("1.8M", "Scans Processed", "since Jan 2025 pilot\n(42% MoM growth)"),
        ("₹4.7 Cr", "ARR", "as of March 2026\n(NRR 128%)"),
    ]
    for i, (val, label, sub) in enumerate(metrics_row1):
        x = 0.8 + i * 2.9
        add_rounded_rect(slide, x, 1.4, 2.6, 1.8)
        add_textbox(slide, x, 1.55, 2.6, 0.5, val, 28,
                    ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, 2.05, 2.6, 0.35, label, 12,
                    WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x, 2.4, 2.6, 0.7, sub, 9,
                    LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Timeline milestones
    milestones = [
        ("Q2 2024", "Founded. Secured ₹2.5 Cr pre-seed from India Quotient & Titan Capital."),
        ("Q4 2024", "Completed clinical validation at AIIMS Jodhpur — 96.3% concordance with senior radiologists (n=12,400)."),
        ("Q1 2025", "Signed MoU with Rajasthan NHM (National Health Mission) — 150 PHCs in first batch."),
        ("Q3 2025", "Expanded to UP & MP. Partnership with SRL Diagnostics for lab-machine integration."),
        ("Q1 2026", "Crossed ₹4 Cr ARR. Selected for WHO Digital Health Innovation Cohort. FDA 510(k) pre-submission filed."),
    ]
    for i, (date, text) in enumerate(milestones):
        y = 3.5 + i * 0.55
        add_textbox(slide, 0.8, y, 1.4, 0.5, date, 10,
                    ACCENT_TEAL, bold=True)
        add_textbox(slide, 2.3, y, 6.8, 0.5, text, 10,
                    LIGHT_GRAY)


def slide_06_business_model(prs):
    """Slide 6 – Business Model & Revenue"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "BUSINESS MODEL", 28,
                ACCENT_BLUE, bold=True)
    add_accent_line(slide, 0.8, 1.05, 1.2, ACCENT_BLUE)

    # Revenue streams
    streams = [
        ("B2G SaaS License", "65% of Revenue",
         "Annual per-PHC license (₹28,000/yr) sold through state NHM tenders. "
         "3-year lock-in with auto-renewal. Includes OTA model updates & support."),
        ("B2B API Access", "20% of Revenue",
         "Diagnostic chains (SRL, Thyrocare, Metropolis) pay ₹3.50 per scan analysis "
         "via API. Minimum commitment: 50,000 scans/month."),
        ("Hardware Bundle", "10% of Revenue",
         "Pre-configured NovaMed Edge Kit (tablet + AI module + adapters) at ₹45,000. "
         "30% gross margin. Sold to PHCs without existing tablets."),
        ("Data Insights (De-identified)", "5% of Revenue",
         "Aggregated, anonymized epidemiological data sold to pharma companies "
         "and public health researchers. HIPAA & DPDP Act compliant."),
    ]

    for i, (name, pct, desc) in enumerate(streams):
        y = 1.4 + i * 1.25
        add_rounded_rect(slide, 0.8, y, 8.2, 1.1)
        add_textbox(slide, 1.0, y + 0.08, 4.5, 0.3, name, 13,
                    WHITE, bold=True)
        add_textbox(slide, 6.5, y + 0.08, 2.3, 0.3, pct, 12,
                    ACCENT_TEAL, bold=True, alignment=PP_ALIGN.RIGHT)
        add_textbox(slide, 1.0, y + 0.42, 7.8, 0.65, desc, 10,
                    LIGHT_GRAY)

    add_multi_text(slide, 0.8, 6.0, 8.0, 1.0, [
        {"text": "Unit Economics (per PHC):", "size": 12, "color": ACCENT_BLUE, "bold": True, "space_after": 4},
        {"text": "LTV: ₹84,000 (3-yr avg tenure)  |  CAC: ₹6,200 (govt tender channel)  |  LTV:CAC = 13.5x  |  Payback: 2.6 months",
         "size": 10, "color": WHITE, "bold": False},
    ])


def slide_07_competition(prs):
    """Slide 7 – Competitive Landscape"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "COMPETITIVE LANDSCAPE", 28,
                ACCENT_TEAL, bold=True)
    add_accent_line(slide, 0.8, 1.05, 1.2, ACCENT_TEAL)

    # Table header
    headers = ["Feature", "NovaMed AI", "Qure.ai", "Google Health", "5C Network"]
    widths  = [1.8, 1.6, 1.6, 1.6, 1.6]
    x_start = 0.8
    y = 1.4
    for i, (h, w) in enumerate(zip(headers, widths)):
        x = x_start + sum(widths[:i])
        add_rounded_rect(slide, x, y, w - 0.05, 0.4,
                         fill_color=ACCENT_BLUE if i == 1 else MEDIUM_GRAY)
        add_textbox(slide, x + 0.05, y + 0.05, w - 0.15, 0.3, h, 10,
                    WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Table rows
    rows = [
        ["Offline / Edge AI",         "✅  Full offline",  "❌  Cloud only",  "❌  Cloud only",  "⚠️  Partial"],
        ["Radiology + Pathology",     "✅  Both",          "✅  Radiology",   "✅  Radiology",   "✅  Radiology"],
        ["Multilingual Voice UX",     "✅  5 languages",   "❌  English",     "❌  English",     "❌  English"],
        ["Govt. Tender Ready",        "✅  GEM listed",    "⚠️  Limited",    "❌  No",          "⚠️  Limited"],
        ["Price (per PHC / yr)",      "₹28,000",           "₹1,20,000",      "N/A (pilot)",     "₹85,000"],
        ["Clinical Validation (India)","AIIMS, CMC",       "Narayana",        "Apollo (US)",     "Manipal"],
        ["Edge Hardware Bundle",      "✅  ₹45K kit",     "❌",              "❌",              "❌"],
    ]

    for r_idx, row in enumerate(rows):
        ry = 1.9 + r_idx * 0.52
        bg = CARD_BG if r_idx % 2 == 0 else MEDIUM_GRAY
        for c_idx, (cell, w) in enumerate(zip(row, widths)):
            x = x_start + sum(widths[:c_idx])
            cell_bg = RGBColor(0x0A, 0x2A, 0x40) if c_idx == 1 and r_idx % 2 == 0 else \
                      RGBColor(0x0D, 0x30, 0x4A) if c_idx == 1 else bg
            add_rounded_rect(slide, x, ry, w - 0.05, 0.45, fill_color=cell_bg)
            col = ACCENT_TEAL if c_idx == 1 else LIGHT_GRAY
            add_textbox(slide, x + 0.05, ry + 0.05, w - 0.15, 0.35, cell, 8,
                        col, bold=(c_idx == 1), alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 0.8, 5.7, 8.0, 0.6,
                "Our moat: Only player with offline-first edge AI + multilingual voice + "
                "government procurement readiness — purpose-built for Bharat's PHCs.",
                12, ORANGE, bold=True)


def slide_08_team(prs):
    """Slide 8 – Team"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "THE TEAM", 28,
                ACCENT_BLUE, bold=True)
    add_accent_line(slide, 0.8, 1.05, 1.2, ACCENT_BLUE)

    team = [
        ("Dr. Ananya Sharma", "CEO & Co-Founder",
         "MBBS (AIIMS Delhi), MBA (ISB Hyderabad). "
         "Former McKinsey Healthcare Practice (3 yrs). "
         "Led digital health rollout for Rajasthan NHM reaching 2M patients. "
         "Published 4 papers on AI diagnostics in The Lancet Digital Health."),
        ("Vikram Patel", "CTO & Co-Founder",
         "B.Tech CS (IIT Bombay), MS ML (Stanford). "
         "Ex-Google Brain — shipped on-device ML for Android Health Sensors. "
         "3 patents in edge inference optimization. "
         "Built Qualcomm's first medical-grade AI benchmark suite."),
        ("Dr. Priya Menon", "Chief Medical Officer",
         "MD Radiology (CMC Vellore), Fellowship (Johns Hopkins). "
         "15 years in tele-radiology, interpreted 200K+ scans. "
         "Board member of Indian Radiology & Imaging Association. "
         "Designed NovaMed's clinical validation protocol."),
        ("Rajesh Iyer", "VP of Government Relations",
         "IAS (Retd.), former Additional Secretary, MoHFW. "
         "Architected the Ayushman Bharat hospital empanelment process. "
         "Personal network across 18 state NHM directors. "
         "Led ₹4,000 Cr procurement for national health IT infrastructure."),
    ]

    for i, (name, role, bio) in enumerate(team):
        y = 1.3 + i * 1.35
        add_rounded_rect(slide, 0.8, y, 8.2, 1.2)
        # Circle placeholder for headshot
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1.0), Inches(y + 0.15), Inches(0.9), Inches(0.9)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT_BLUE
        circle.line.fill.background()
        initials = "".join(w[0] for w in name.split() if w[0].isupper())
        # Add initials text
        tf = circle.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = initials
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(8)

        add_textbox(slide, 2.2, y + 0.1, 3.0, 0.3, name, 14,
                    WHITE, bold=True)
        add_textbox(slide, 2.2, y + 0.4, 3.0, 0.25, role, 11,
                    ACCENT_TEAL, bold=True)
        add_textbox(slide, 2.2, y + 0.65, 6.6, 0.55, bio, 9,
                    LIGHT_GRAY)


def slide_09_financials(prs):
    """Slide 9 – Financials & Use of Funds"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    add_textbox(slide, 0.8, 0.5, 8.0, 0.6, "FINANCIALS & USE OF FUNDS", 28,
                ACCENT_TEAL, bold=True)
    add_accent_line(slide, 0.8, 1.05, 1.2, ACCENT_TEAL)

    # Financial projections table
    add_rounded_rect(slide, 0.8, 1.3, 8.2, 2.4)

    headers_fin = ["Metric", "FY25 (Actual)", "FY26 (Current)", "FY27 (Proj)", "FY28 (Proj)"]
    widths_fin = [2.0, 1.5, 1.5, 1.5, 1.5]
    for i, (h, w) in enumerate(zip(headers_fin, widths_fin)):
        x = 0.9 + sum(widths_fin[:i])
        add_textbox(slide, x, 1.4, w, 0.3, h, 10, ACCENT_BLUE, bold=True,
                    alignment=PP_ALIGN.CENTER)

    fin_rows = [
        ["Revenue", "₹1.2 Cr", "₹4.7 Cr", "₹18.5 Cr", "₹52 Cr"],
        ["Gross Margin", "58%", "64%", "72%", "78%"],
        ["PHCs Active", "85", "312", "1,200", "3,800"],
        ["Burn Rate (mo)", "₹38L", "₹52L", "₹65L", "₹48L"],
        ["Runway", "8 months", "Post-raise: 24 mo", "—", "Cash-flow +ve"],
    ]
    for r_idx, row in enumerate(fin_rows):
        ry = 1.75 + r_idx * 0.38
        for c_idx, (cell, w) in enumerate(zip(row, widths_fin)):
            x = 0.9 + sum(widths_fin[:c_idx])
            c = WHITE if c_idx == 0 else LIGHT_GRAY
            b = c_idx == 0
            add_textbox(slide, x, ry, w, 0.3, cell, 10, c, bold=b,
                        alignment=PP_ALIGN.CENTER)

    # Use of Funds
    add_textbox(slide, 0.8, 3.9, 4.0, 0.4, "Raising: $12M Series A (₹100 Cr)", 14,
                WHITE, bold=True)
    add_textbox(slide, 0.8, 4.25, 4.0, 0.3, "Lead: Nexus Venture Partners (term sheet signed)", 10,
                ACCENT_TEAL)

    funds = [
        ("R&D & Product (40%)", "₹40 Cr", "Pathology V2 engine, FDA clearance, multilingual expansion to 12 languages"),
        ("GTM & Expansion (30%)", "₹30 Cr", "Scale to 8 states, build 15-person govt sales team, attend 24 state NHM tenders"),
        ("Infrastructure (15%)", "₹15 Cr", "Edge hardware inventory, HIPAA-grade cloud infra, SOC 2 Type II certification"),
        ("Team & Ops (15%)", "₹15 Cr", "Hire 60 engineers (ML, embedded, full-stack), open Bengaluru R&D center"),
    ]
    for i, (name, amt, desc) in enumerate(funds):
        y = 4.7 + i * 0.6
        add_rounded_rect(slide, 0.8, y, 8.2, 0.55)
        add_textbox(slide, 1.0, y + 0.05, 2.8, 0.25, name, 10,
                    ACCENT_BLUE, bold=True)
        add_textbox(slide, 3.8, y + 0.05, 1.0, 0.25, amt, 10,
                    ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
        add_textbox(slide, 4.9, y + 0.05, 3.9, 0.45, desc, 9,
                    LIGHT_GRAY)


def slide_10_ask(prs):
    """Slide 10 – The Ask & Contact"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_BG)

    add_accent_line(slide, 1.5, 1.5, 2.0, ACCENT_TEAL)

    add_textbox(slide, 1.5, 1.7, 7.0, 0.8, "THE ASK", 36,
                ACCENT_BLUE, bold=True, alignment=PP_ALIGN.LEFT)

    add_multi_text(slide, 1.5, 2.6, 7.0, 2.0, [
        {"text": "$12M Series A to scale AI diagnostics", "size": 20, "color": WHITE, "bold": True, "space_after": 6},
        {"text": "from 312 to 3,800 PHCs by FY28", "size": 20, "color": ACCENT_TEAL, "bold": True, "space_after": 20},
        {"text": "With your partnership, we will:", "size": 13, "color": LIGHT_GRAY, "space_after": 8},
        {"text": "✦  Save an estimated 42,000 lives annually through early detection", "size": 12, "color": WHITE, "space_after": 5},
        {"text": "✦  Reduce diagnostic costs by 85% for India's most underserved communities", "size": 12, "color": WHITE, "space_after": 5},
        {"text": "✦  Build the largest rural health AI dataset globally (50M+ data points by 2028)", "size": 12, "color": WHITE, "space_after": 5},
        {"text": "✦  Achieve cash-flow positivity by Q3 FY28 with 78% gross margins", "size": 12, "color": WHITE},
    ])

    add_rounded_rect(slide, 1.5, 5.2, 7.0, 1.5)
    add_multi_text(slide, 1.8, 5.3, 6.5, 1.3, [
        {"text": "Let's talk.", "size": 18, "color": ACCENT_BLUE, "bold": True, "space_after": 10},
        {"text": "Dr. Ananya Sharma  —  CEO & Co-Founder", "size": 12, "color": WHITE, "bold": True, "space_after": 4},
        {"text": "📧  ananya@novamed.ai   |   📞  +91-98765-43210", "size": 11, "color": LIGHT_GRAY, "space_after": 3},
        {"text": "🌐  www.novamed.ai   |   LinkedIn: /in/ananya-sharma-novamed", "size": 11, "color": LIGHT_GRAY},
    ])


# ═══════════════════════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_market(prs)
    slide_05_traction(prs)
    slide_06_business_model(prs)
    slide_07_competition(prs)
    slide_08_team(prs)
    slide_09_financials(prs)
    slide_10_ask(prs)

    os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"[OK] Pitch deck saved to: {OUTPUT_PATH}")
    print(f"   Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
