import os
import pdfplumber
from pptx import Presentation

# Try importing OCR dependencies
try:
    import fitz  # PyMuPDF
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    from PIL import Image
    import pytesseract
    import io
    HAS_OCR = True
except ImportError:
    HAS_OCR = False


def _ocr_page_with_fitz(pdf_path, page_number):
    """Use PyMuPDF to extract text from a specific page (handles embedded text + images better)."""
    if not HAS_FITZ:
        return ""
    try:
        doc = fitz.open(pdf_path)
        page = doc[page_number]
        text = page.get_text("text").strip()
        
        # If still no text, try extracting text from images on the page via OCR
        if not text and HAS_OCR:
            text = _ocr_page_images(page)
        
        doc.close()
        return text
    except Exception as e:
        print(f"  PyMuPDF fallback failed for page {page_number + 1}: {e}")
        return ""


def _ocr_page_images(page):
    """Extract images from a PyMuPDF page and run OCR on them."""
    if not HAS_OCR:
        return ""
    try:
        texts = []
        image_list = page.get_images(full=True)
        for img_info in image_list:
            xref = img_info[0]
            doc = page.parent
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            
            image = Image.open(io.BytesIO(image_bytes))
            ocr_text = pytesseract.image_to_string(image).strip()
            if ocr_text:
                texts.append(ocr_text)
        return " ".join(texts)
    except Exception as e:
        return ""


def extract_from_pdf(pdf_path):
    """Extracts text from a PDF pitch deck slide by slide.
    
    Uses a multi-tier extraction strategy:
    1. pdfplumber (fast, handles selectable text)
    2. PyMuPDF/fitz (handles more PDF formats and embedded text)
    3. pytesseract OCR (handles scanned/image-only pages)
    """
    slides_data = []
    empty_count = 0
    
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for i, page in enumerate(pdf.pages):
                text = page.extract_text()
                if text:
                    text = text.strip()
                
                # Fallback: if pdfplumber got nothing, try PyMuPDF
                if not text:
                    text = _ocr_page_with_fitz(pdf_path, i)
                    
                if not text:
                    empty_count += 1
                    
                slides_data.append({
                    "slide_number": i + 1,
                    "raw_text": text or ""
                })
    except Exception as e:
        print(f"Error reading PDF {pdf_path}: {e}")
    
    total = len(slides_data)
    if empty_count > 0:
        print(f"  Note: {empty_count}/{total} pages had no extractable text (image-based slides).")
        if empty_count == total and not HAS_FITZ:
            print("  Tip: Install PyMuPDF for better extraction: pip install PyMuPDF")
        if empty_count == total and not HAS_OCR:
            print("  Tip: Install pytesseract + Pillow for OCR: pip install pytesseract Pillow")
    
    return slides_data

def extract_from_pptx(pptx_path):
    """Extracts text from a PPTX pitch deck slide by slide."""
    slides_data = []
    try:
        prs = Presentation(pptx_path)
        for i, slide in enumerate(prs.slides):
            text_runs = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
            text = " ".join(text_runs).strip()
            slides_data.append({
                "slide_number": i + 1,
                "raw_text": text
            })
    except Exception as e:
        print(f"Error reading PPTX {pptx_path}: {e}")
    return slides_data

def extract_text(file_path):
    """Router to handle both PDF and PPTX extraction based on extension."""
    _, ext = os.path.splitext(file_path)
    if ext.lower() == '.pdf':
        return extract_from_pdf(file_path)
    elif ext.lower() == '.pptx':
        return extract_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}. Only .pdf and .pptx are supported.")
